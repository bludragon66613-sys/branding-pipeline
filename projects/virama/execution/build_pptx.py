"""
Virāma by SSquare — Brand Presentation PPTX Builder
16-slide deck | Brand palette | Cormorant Garamond + DM Sans
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Dimensions ─────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Brand Palette ───────────────────────────────────────────────────
STONE       = RGBColor(0x2C, 0x28, 0x25)   # Virāma Stone
WHITE       = RGBColor(0xF5, 0xF1, 0xEC)   # Pavilion White
GOLD        = RGBColor(0xC9, 0xA9, 0x6E)   # Arcadian Gold
SAND        = RGBColor(0xE8, 0xE2, 0xD8)   # Sandstone Pale
GREEN       = RGBColor(0x3D, 0x4A, 0x3A)   # The Green Circuit
STONE_SOFT  = RGBColor(0x4A, 0x44, 0x40)   # Soft stone for body

# ── Font names ───────────────────────────────────────────────────────
SERIF   = "Cormorant Garamond"
SANS    = "DM Sans"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def rect(slide, x, y, w, h, fill_rgb=None, fill_alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    return shape

def txbox(slide, text, x, y, w, h,
          font_name=SANS, size=12, bold=False, italic=False,
          color=WHITE, align=PP_ALIGN.LEFT,
          line_spacing=None, space_before=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def multiline_txbox(slide, lines, x, y, w, h,
                    font_name=SERIF, size=40, italic=True,
                    color=WHITE, align=PP_ALIGN.LEFT,
                    line_gap=8, bold=False):
    """Each item in lines is a string; renders each on its own paragraph."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(line_gap)
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.italic = italic
        run.font.bold = bold
        run.font.color.rgb = color
    return tb

def hairline(slide, x, y, length, color=GOLD, width_pt=0.75):
    line = slide.shapes.add_connector(1, x, y, x + length, y)
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line

def eyebrow(slide, text, x, y, w):
    """Tracked small caps eyebrow in gold"""
    tb = txbox(slide, text, x, y, w, Inches(0.25),
               font_name=SANS, size=8, color=GOLD, bold=False)
    return tb

def dark_slide(slide):
    rect(slide, 0, 0, W, H, fill_rgb=STONE)

def sand_slide(slide):
    rect(slide, 0, 0, W, H, fill_rgb=SAND)

def white_slide(slide):
    rect(slide, 0, 0, W, H, fill_rgb=WHITE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 01 — Cover
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

# Centred vertical layout
multiline_txbox(s, ["Virāma"],
    Inches(1.5), Inches(1.8), Inches(10), Inches(1.8),
    font_name=SERIF, size=80, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "विराम",
    Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
    font_name=SERIF, size=28, italic=False, color=SAND,
    align=PP_ALIGN.CENTER)

hairline(s, Inches(5.2), Inches(4.35), Inches(2.9))

txbox(s, "TWELVE PRIVATE VILLAS  ·  JUBILEE HILLS  ·  HYDERABAD",
    Inches(1.5), Inches(4.55), Inches(10), Inches(0.4),
    font_name=SANS, size=8, color=GOLD, align=PP_ALIGN.CENTER)

txbox(s, "A SSquare Development",
    Inches(1.5), Inches(6.7), Inches(10), Inches(0.4),
    font_name=SANS, size=9, italic=True, color=SAND, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 02 — The Pause
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

# Right panel — lighter tone to suggest split
rp = rect(s, Inches(5.5), 0, Inches(7.83), H, fill_rgb=RGBColor(0x38,0x33,0x2F))

M = Inches(1.2)
eyebrow(s, "VIRĀMA  ·  विराम", M, Inches(2.0), Inches(4.5))
hairline(s, M, Inches(2.35), Inches(2.8))

multiline_txbox(s, ["The word does not mean stop.", "It means the breath between."],
    M, Inches(2.55), Inches(4.2), Inches(2.0),
    font_name=SERIF, size=36, italic=True, color=WHITE, line_gap=6)

txbox(s, "In Sanskrit, virāma is the pause that gives a phrase its meaning. Not silence as absence — silence as presence. We named this place after that idea. Twelve homes built not around what fills them, but around what they let go of.",
    M, Inches(4.8), Inches(4.2), Inches(1.8),
    font_name=SANS, size=10.5, color=SAND)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 03 — SSquare 35 Years
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
white_slide(s)

M = Inches(1.2)
eyebrow(s, "THE DEVELOPER", M, Inches(1.4), Inches(4))
hairline(s, M, Inches(1.75), Inches(2.5), color=GOLD)

multiline_txbox(s, ["Thirty-five years of building", "things that last."],
    M, Inches(1.95), Inches(7), Inches(1.8),
    font_name=SERIF, size=40, italic=True, color=STONE)

txbox(s, "SSquare was founded in 1991. Since then, we have developed over ten million square feet of premium commercial real estate across Hyderabad. We have never built for volume. We have built to a standard. Virāma is our first residential work — and it will be our definitive one.",
    M, Inches(3.95), Inches(7), Inches(1.6),
    font_name=SANS, size=11, color=STONE_SOFT)

hairline(s, M, Inches(5.75), Inches(3), color=GOLD)
txbox(s, '"The Standard, Privately Held."',
    M, Inches(5.95), Inches(8), Inches(0.8),
    font_name=SERIF, size=22, italic=True, color=GOLD)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 04 — Jubilee Hills / Malibu
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

M = Inches(1.2)
eyebrow(s, "THE LOCATION", M, Inches(1.8), Inches(4))
hairline(s, M, Inches(2.15), Inches(2.5))

multiline_txbox(s, ["Where Hyderabad's hills", "meet their quietest road."],
    M, Inches(2.35), Inches(5.5), Inches(1.8),
    font_name=SERIF, size=38, italic=True, color=WHITE)

txbox(s, "Prashashan Nagar. A single street in Jubilee Hills known to its residents simply as Malibu. Gated, green, unhurried — 500 metres from Film Nagar, four minutes from the Jubilee Hills Club, and entirely apart from everything else. The address is not a postcode. It is a posture.",
    M, Inches(4.35), Inches(5.5), Inches(1.8),
    font_name=SANS, size=10.5, color=SAND)

# Distance callouts — right side
dist_x = Inches(8.5)
dists = [
    ("4 min", "Jubilee Hills Club"),
    ("10 min", "HITECH City corridor"),
    ("10 min", "KBR National Park"),
    ("18 min", "Financial District"),
]
for i, (time, dest) in enumerate(dists):
    y_pos = Inches(1.8 + i * 1.2)
    txbox(s, time, dist_x, y_pos, Inches(1.5), Inches(0.45),
          font_name=SERIF, size=22, italic=True, color=GOLD, align=PP_ALIGN.LEFT)
    txbox(s, dest, dist_x, y_pos + Inches(0.38), Inches(4), Inches(0.35),
          font_name=SANS, size=9, color=SAND)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 05 — The Estate
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
sand_slide(s)

M = Inches(1.2)
eyebrow(s, "THE ESTATE", M, Inches(1.4), Inches(4))
hairline(s, M, Inches(1.75), Inches(2.5), color=GOLD)

txbox(s, "Twelve. No more.",
    M, Inches(1.95), Inches(5), Inches(1.0),
    font_name=SERIF, size=52, italic=True, color=STONE)

txbox(s, "A gated estate of twelve private villas. Each home has its own pool, its own forecourt, its own silence. The Green Circuit — a continuous landscape corridor of native planting — connects them without crowding them. This is not a development. It is a private enclave.",
    M, Inches(3.3), Inches(5.2), Inches(1.8),
    font_name=SANS, size=11, color=STONE_SOFT)

# Stats block — right side
stats = [("12", "Private Villas"), ("1", "Gate · One Entry"), ("400m", "The Green Circuit")]
for i, (num, label) in enumerate(stats):
    sx = Inches(8.0 + i * 1.7)
    txbox(s, num, sx, Inches(2.8), Inches(1.5), Inches(1.2),
          font_name=SERIF, size=52, italic=False, color=STONE, align=PP_ALIGN.CENTER)
    txbox(s, label, sx, Inches(3.9), Inches(1.5), Inches(0.5),
          font_name=SANS, size=8, color=STONE_SOFT, align=PP_ALIGN.CENTER)
    hairline(s, sx + Inches(0.2), Inches(4.6), Inches(1.1), color=GOLD)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 06 — Architecture Language
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

M = Inches(1.2)
eyebrow(s, "ARCHITECTURE", M, Inches(1.6), Inches(4))
hairline(s, M, Inches(1.95), Inches(2.5))

multiline_txbox(s, ["Stone, glass, and the", "discipline of restraint."],
    M, Inches(2.15), Inches(5.2), Inches(1.8),
    font_name=SERIF, size=38, italic=True, color=WHITE)

txbox(s, "Single storey. Honed limestone walls. Dark basalt at the base. A cantilevered roof that extends the interior into the landscape. Floor-to-ceiling glazing that disappears at dusk. A teak pivot door that announces arrival without announcing itself. Every material chosen for how it weathers.",
    M, Inches(4.2), Inches(5.2), Inches(1.9),
    font_name=SANS, size=10.5, color=SAND)

# Material callouts — right column
materials = ["Honed limestone", "Dark basalt base", "Cantilevered roof", "Teak pivot door", "Floor-to-ceiling glazing"]
mat_x = Inches(8.8)
for i, mat in enumerate(materials):
    rect(s, mat_x, Inches(1.5 + i * 1.05), Inches(0.04), Inches(0.6), fill_rgb=GOLD)
    txbox(s, mat, mat_x + Inches(0.2), Inches(1.5 + i * 1.05), Inches(3.8), Inches(0.5),
          font_name=SANS, size=10, color=SAND)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 07 — The Villa
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

M = Inches(1.2)
eyebrow(s, "THE VILLA", M, Inches(1.8), Inches(4))
hairline(s, M, Inches(2.15), Inches(2.5))

multiline_txbox(s, ["A home that knows", "when to step back."],
    M, Inches(2.35), Inches(5.2), Inches(1.8),
    font_name=SERIF, size=42, italic=True, color=WHITE)

txbox(s, "The living spaces open fully to the outdoors. There is no front of house and back of house — the whole villa is the front. Bedrooms face the pool. The kitchen faces the garden. The roof is a horizontal plane that makes the sky feel deliberately framed. Living here is an act of attention.",
    M, Inches(4.4), Inches(5.2), Inches(2.0),
    font_name=SANS, size=10.5, color=SAND)

# Right — image placeholder box with label
rbox = rect(s, Inches(7.5), Inches(0.5), Inches(5.6), Inches(6.5),
            fill_rgb=RGBColor(0x38, 0x33, 0x2F))
txbox(s, "[ Villa interior photography ]",
    Inches(7.6), Inches(3.4), Inches(5.4), Inches(0.5),
    font_name=SANS, size=9, color=RGBColor(0x60,0x58,0x52), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 08 — Private Pool
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

# Full-bleed image placeholder
rbox = rect(s, 0, 0, W, H, fill_rgb=RGBColor(0x1A, 0x2A, 0x2F))
txbox(s, "[ Private pool photography — dawn, still water, limestone deck ]",
    Inches(3), Inches(3.3), Inches(7.33), Inches(0.6),
    font_name=SANS, size=9, color=RGBColor(0x55,0x65,0x60), align=PP_ALIGN.CENTER)

# Bottom-left type
eyebrow(s, "PRIVATE POOL", Inches(1.2), Inches(5.6), Inches(4))
hairline(s, Inches(1.2), Inches(5.95), Inches(3.0))
txbox(s, "Yours. Only yours.",
    Inches(1.2), Inches(6.1), Inches(6), Inches(0.9),
    font_name=SERIF, size=48, italic=True, color=WHITE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 09 — The Green Circuit
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
# Dark green background
rect(s, 0, 0, W, H, fill_rgb=RGBColor(0x22, 0x2C, 0x22))

M = Inches(1.2)
eyebrow(s, "THE GREEN CIRCUIT", M, Inches(1.6), Inches(5))
hairline(s, M, Inches(1.95), Inches(2.8))

multiline_txbox(s, ["Native trees.", "Continuous ground.", "No fences between."],
    M, Inches(2.15), Inches(5.2), Inches(2.4),
    font_name=SERIF, size=38, italic=True, color=WHITE, line_gap=4)

txbox(s, "The landscape at Virāma is not decoration. The Green Circuit is a continuous corridor of native Deccan plantings — fig, neem, Indian coral — that weaves between and behind the villas without interruption. The estate breathes as one organism. You are not separated from your neighbours by walls. You are separated by garden.",
    M, Inches(4.85), Inches(5.4), Inches(1.9),
    font_name=SANS, size=10.5, color=RGBColor(0xC8,0xD4,0xC0))

# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — Arrival
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

# Left dark panel
rect(s, 0, 0, Inches(5.2), H, fill_rgb=STONE)
# Right image area
rbox = rect(s, Inches(4.8), 0, Inches(8.53), H, fill_rgb=RGBColor(0x30,0x28,0x20))
txbox(s, "[ Teak pivot door — straight-on, door-handle height ]",
    Inches(5.0), Inches(3.3), Inches(8.1), Inches(0.5),
    font_name=SANS, size=9, color=RGBColor(0x60,0x50,0x40), align=PP_ALIGN.CENTER)

M = Inches(1.0)
eyebrow(s, "ARRIVAL", M, Inches(1.8), Inches(3.8))
hairline(s, M, Inches(2.15), Inches(2.5))

multiline_txbox(s, ["The door is the first", "thing you understand."],
    M, Inches(2.35), Inches(3.8), Inches(1.8),
    font_name=SERIF, size=36, italic=True, color=WHITE)

txbox(s, "A teak pivot door, floor-to-ceiling, set in a deep limestone reveal. The forecourt is quiet — no waterfalls, no statement planting. Basalt stepping stones across pale gravel. The arrival sequence is brief and unhurried. You are not announced. You simply come home.",
    M, Inches(4.3), Inches(3.8), Inches(1.9),
    font_name=SANS, size=10.5, color=SAND)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — The Pavilion
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

rbox = rect(s, 0, 0, W, H, fill_rgb=RGBColor(0x24, 0x20, 0x1C))
txbox(s, "[ Common pavilion — wide angle, open sides, garden beyond ]",
    Inches(3), Inches(3.4), Inches(7.33), Inches(0.5),
    font_name=SANS, size=9, color=RGBColor(0x55,0x50,0x48), align=PP_ALIGN.CENTER)

eyebrow(s, "THE PAVILION", Inches(1.2), Inches(5.3), Inches(5))
hairline(s, Inches(1.2), Inches(5.65), Inches(2.8))
multiline_txbox(s, ["A room that belongs", "to all twelve."],
    Inches(1.2), Inches(5.82), Inches(6), Inches(1.2),
    font_name=SERIF, size=38, italic=True, color=WHITE, line_gap=4)

txbox(s, "The Pavilion is the estate's single shared space — open on three sides to the garden. No gym, no concierge desk, no co-working nook. There is a long table, good light, and the sound of wind through the canopy.",
    Inches(1.2), Inches(4.0), Inches(5.5), Inches(1.2),
    font_name=SANS, size=10.5, color=SAND)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Facts
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
sand_slide(s)

M = Inches(1.2)
eyebrow(s, "AT A GLANCE", M, Inches(0.8), Inches(4))
hairline(s, M, Inches(1.15), Inches(2.5), color=GOLD)

multiline_txbox(s, ["Twelve villas.", "One estate.", "One address."],
    M, Inches(1.35), Inches(5), Inches(2.0),
    font_name=SERIF, size=36, italic=True, color=STONE, line_gap=4)

facts_left = [
    ("VILLAS", "12 private homes, each with individual pool"),
    ("PRICE", "Rs 13 Crore to Rs 20 Crore"),
    ("LOCATION", "Prashashan Nagar, Jubilee Hills, Hyderabad"),
    ("CONFIGURATION", "Single storey · 4 BHK + Study"),
]
facts_right = [
    ("POSSESSION", "[Possession quarter / year — TBC]"),
    ("RERA", "[RERA registration number — TBC]"),
    ("DEVELOPER", "SSquare · Est. 1991"),
    ("TAGLINE", '"The Standard, Privately Held."'),
]

def fact_block(slide, facts, x, y_start):
    for i, (label, val) in enumerate(facts):
        y = y_start + i * Inches(1.1)
        txbox(slide, label, x, y, Inches(3.5), Inches(0.3),
              font_name=SANS, size=7.5, color=STONE_SOFT, bold=False)
        hairline(slide, x, y + Inches(0.28), Inches(3.2), color=GOLD, width_pt=0.5)
        txbox(slide, val, x, y + Inches(0.35), Inches(3.8), Inches(0.55),
              font_name=SANS, size=11, color=STONE)

fact_block(s, facts_left,  Inches(1.2), Inches(3.6))
fact_block(s, facts_right, Inches(7.0), Inches(3.6))

# vertical divider
hairline(s, Inches(6.65), Inches(3.6), Inches(0), color=GOLD)
# Actually a vertical line requires different approach — skip for now

# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — Location Map
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
sand_slide(s)

M = Inches(1.2)
eyebrow(s, "LOCATION CONNECTIVITY", M, Inches(0.9), Inches(5))
hairline(s, M, Inches(1.25), Inches(2.5), color=GOLD)

multiline_txbox(s, ["The city at a distance.", "Everything you need", "within four minutes."],
    M, Inches(1.45), Inches(4.5), Inches(1.8),
    font_name=SERIF, size=32, italic=True, color=STONE, line_gap=4)

connections = [
    ("4 min",  "Jubilee Hills Club"),
    ("6 min",  "Jubilee Hills Check Post"),
    ("8 min",  "Banjara Hills Road No. 12"),
    ("10 min", "KBR National Park"),
    ("12 min", "Hyderabad International Convention Centre"),
    ("18 min", "Financial District · HITECH City"),
]
for i, (time, dest) in enumerate(connections):
    y = Inches(3.5 + i * 0.6)
    txbox(s, time, M, y, Inches(1.4), Inches(0.45),
          font_name=SERIF, size=18, italic=True, color=GOLD)
    txbox(s, dest, Inches(2.7), y + Inches(0.05), Inches(3.8), Inches(0.4),
          font_name=SANS, size=10, color=STONE_SOFT)

# Right — illustrated map placeholder
rbox = rect(s, Inches(7.2), Inches(0.5), Inches(5.9), Inches(6.5),
            fill_rgb=RGBColor(0xD8, 0xD2, 0xC8))
txbox(s, "[ Illustrated estate map — fine-line, Arcadian Gold on Sandstone ]",
    Inches(7.3), Inches(3.5), Inches(5.7), Inches(0.5),
    font_name=SANS, size=8.5, color=STONE_SOFT, align=PP_ALIGN.CENTER)
# Dot for location
rect(s, Inches(10.1), Inches(3.1), Inches(0.15), Inches(0.15), fill_rgb=GOLD)
txbox(s, "Virāma", Inches(10.3), Inches(3.05), Inches(1.5), Inches(0.35),
      font_name=SERIF, size=11, italic=True, color=STONE)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 14 — SSquare Track Record
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
white_slide(s)

M = Inches(1.2)
eyebrow(s, "SSQUARE  ·  EST. 1991", M, Inches(0.9), Inches(5))
hairline(s, M, Inches(1.25), Inches(2.5), color=GOLD)

multiline_txbox(s, ["The people who built Hyderabad's skyline", "now build its quietest home."],
    M, Inches(1.45), Inches(11), Inches(1.6),
    font_name=SERIF, size=34, italic=True, color=STONE)

txbox(s, "Over three decades and ten million square feet of premium commercial development in Hyderabad. Grade A offices, institutional campuses, long-term assets. SSquare has never listed, never diluted, never rushed. The same discipline that governs a 500,000 sq ft commercial campus governs a single Virāma villa.",
    M, Inches(3.25), Inches(11), Inches(1.4),
    font_name=SANS, size=11, color=STONE_SOFT)

# Triptych stats
stats3 = [("35", "Years in Hyderabad real estate"), ("10M+", "Square feet developed"), ("1", "Residential debut")]
for i, (num, label) in enumerate(stats3):
    sx = Inches(1.4 + i * 4.0)
    txbox(s, num, sx, Inches(4.8), Inches(3.5), Inches(1.4),
          font_name=SERIF, size=72, italic=False, color=STONE, align=PP_ALIGN.CENTER)
    txbox(s, label, sx, Inches(6.1), Inches(3.5), Inches(0.5),
          font_name=SANS, size=9, color=STONE_SOFT, align=PP_ALIGN.CENTER)

hairline(s, Inches(9.8), Inches(5.6), Inches(2.3), color=GOLD)
txbox(s, '"The Standard, Privately Held."',
    Inches(9.8), Inches(5.8), Inches(3.3), Inches(0.6),
    font_name=SERIF, size=13, italic=True, color=GOLD)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 15 — Come Home to Quiet
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

# Full-bleed cinematic image placeholder
rbox = rect(s, 0, 0, W, H, fill_rgb=RGBColor(0x1A, 0x1C, 0x18))
txbox(s, "[ Cinematic dusk photograph — villa from distance, pool reflecting sky, canopy framing sides ]",
    Inches(2), Inches(3.4), Inches(9.33), Inches(0.5),
    font_name=SANS, size=9, color=RGBColor(0x45,0x48,0x42), align=PP_ALIGN.CENTER)

# Centred headline
txbox(s, "Come home to quiet.",
    Inches(1.5), Inches(2.5), Inches(10.33), Inches(1.4),
    font_name=SERIF, size=64, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

hairline(s, Inches(5.2), Inches(4.05), Inches(2.9))

txbox(s, "Virāma by SSquare",
    Inches(1.5), Inches(4.25), Inches(10.33), Inches(0.5),
    font_name=SANS, size=12, color=SAND, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════
# SLIDE 16 — Back Cover
# ═══════════════════════════════════════════════════════════════════
s = add_slide()
dark_slide(s)

txbox(s, "Virāma",
    Inches(1.2), Inches(1.4), Inches(6), Inches(1.4),
    font_name=SERIF, size=64, italic=True, color=WHITE)
txbox(s, "विराम",
    Inches(1.2), Inches(2.7), Inches(4), Inches(0.7),
    font_name=SERIF, size=24, color=SAND)

hairline(s, Inches(1.2), Inches(3.6), Inches(3))

txbox(s, "Sales Enquiries",
    Inches(1.2), Inches(3.9), Inches(5), Inches(0.35),
    font_name=SANS, size=8, color=GOLD)
txbox(s, "[Phone number]  ·  [Email address]  ·  [Website URL]",
    Inches(1.2), Inches(4.28), Inches(7), Inches(0.4),
    font_name=SANS, size=10.5, color=SAND)
txbox(s, "Site Visits by Appointment Only",
    Inches(1.2), Inches(4.85), Inches(6), Inches(0.35),
    font_name=SANS, size=8, color=GOLD)
txbox(s, "Prashashan Nagar, Jubilee Hills, Hyderabad",
    Inches(1.2), Inches(5.23), Inches(7), Inches(0.4),
    font_name=SANS, size=10.5, color=SAND)

hairline(s, Inches(1.2), Inches(6.55), Inches(3))
txbox(s, "A SSquare Development  ·  Est. 1991",
    Inches(1.2), Inches(6.72), Inches(6), Inches(0.4),
    font_name=SANS, size=9, color=STONE_SOFT)
txbox(s, '"The Standard, Privately Held."',
    Inches(1.2), Inches(7.08), Inches(6), Inches(0.4),
    font_name=SERIF, size=11, italic=True, color=GOLD)

# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
out = r"C:\Users\Rohan\branding-pipeline\projects\virama\execution\Virama_SSquare_Brand_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
